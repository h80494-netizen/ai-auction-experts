import os
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, RGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

DEFAULT_DOWNLOADS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "downloads"))

def generate_analysis_doc(data, output_dir=DEFAULT_DOWNLOADS_DIR):
    """
    분석 데이터를 기반으로 심층 분석 Word 문서(.docx)를 생성합니다.
    """
    os.makedirs(output_dir, exist_ok=True)
    
    case_number = data.get("case_number", "알수없음")
    pType = data.get("property_type", "주택")
    house_count = data.get("house_count", 0)
    investor_type = data.get("investor_type", "개인")
    duration = data.get("investment_duration", "단기")
    r_cond = data.get("repair_condition", "부분수리")
    is_reg = data.get("is_regulated_area", False)
    
    doc = Document()
    
    # 스타일 세팅
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Malgun Gothic'
    font.size = Pt(11)
    
    # 0. 커버 페이지 (표지)
    doc.add_heading('AI 기반 부동산 경매 심층 분석 리포트', 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph('\n')
    
    # 썸네일/지도 추가 (crawler가 다운로드한 실제 파일 경로 사용)
    safe_case = case_number.replace(" ", "_").replace("/", "_")
    img_path_photo = os.path.join(output_dir, safe_case, "photo.jpg")
    img_path_map = os.path.join(output_dir, safe_case, "map.jpg")
    img_path_structure = os.path.join(output_dir, safe_case, "structure.jpg")
    
    # (표지에 들어갔던 3장 사진은 권리분석 양식에 맞춰 아래로 이동할 것이므로 삭제)
    p_cover = doc.add_paragraph()
    p_cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cover.add_run(f"사건번호: {case_number}\n").bold = True
    p_cover.add_run(f"상태: {data.get('status', '알수없음')}\n")
    p_cover.add_run(f"주소: {data.get('address', '주소 미상')}\n")
    doc.add_page_break()

    appraised = int(data.get('appraised_value', 3800000000))
    minimum = int(data.get('minimum_value', 3000000000))
    risks = data.get("risks", [])
    real_risks = [r for r in risks if r != "HUG 대항력포기"]
    has_rights = len(real_risks) > 0
    
    # --- 정밀 권리분석 및 베이스라인 타겟가 로직 ---
    address = data.get('address', '')
    base_rate = 0.80
    if '강남구' in address or '서초구' in address or '송파구' in address:
        base_rate = 0.95
    elif '서울' in address:
        base_rate = 0.90
    elif '경기' in address or '인천' in address:
        base_rate = 0.85
        
    if pType == '아파트':
        base_rate += 0.02
    elif pType == '빌라' or '다세대' in pType:
        base_rate -= 0.05
    elif pType == '토지':
        base_rate -= 0.10
        
    target_price = int(appraised * base_rate)
    assumed_debt = 0
    can_invest = True
    force_no_loan = False
    
    # AI 딥리서치가 산출한 가격이 있다면 하드코딩된 지역 기반 비율 대신 사용
    if "ai_sise" in data and "ai_target" in data:
        appraised = data["ai_sise"] # 법원 감정가 대신 최신 AI 시세를 기준으로 삼음
        target_price = data["ai_target"]
        assumed_debt = data.get("ai_insu", 0)
    else:
        # 기존 휴리스틱 로직
        if "법정지상권" in risks:
            if pType == "토지":
                can_invest = False
                target_price = 0
            else:
                target_price = int(appraised * 0.60)
                assumed_debt = int(appraised * 0.20)
                
        if can_invest and "대항력 임차인" in risks:
            if "보증금미상" in risks:
                assumed_debt += int(appraised * 0.60)
                target_price -= assumed_debt
            elif "배당요구" in risks:
                assumed_debt += int(appraised * 0.40)
                target_price -= assumed_debt
            else:
                assumed_debt += int(appraised * 0.35)
                target_price -= assumed_debt

        if can_invest and "유치권" in risks:
            lien_claim = int(appraised * 0.15)
            assumed_debt += lien_claim
            target_price -= lien_claim
        
    if can_invest and "지분매각" in risks:
        target_price = int(target_price * 0.80)
        force_no_loan = True

    if can_invest and "위반건축물" in risks:
        fine = int(appraised * 0.05)
        assumed_debt += fine
        target_price -= fine
        force_no_loan = True

    if can_invest and "미납관리비" in risks:
        assumed_debt += 5000000

    if can_invest and pType == '토지' and any(r in risks for r in ["도로", "하천", "구거"]):
        force_no_loan = True

    # -------------------------------------------------------------
    # (Word 문서 생성 로직) - 권리분석.txt 양식 적용
    # -------------------------------------------------------------
    
    import re
    def extract_section(num, text):
        pattern = r'(?:^|\n)[#\*\s]*' + str(num) + r'\.\s*[^\n]*\n([\s\S]*?)(?=(?:^|\n)[#\*\s]*\d+\.\s*|$)'
        match = re.search(pattern, text, re.IGNORECASE)
        return match.group(1).strip() if match else ""

    if data.get("analysis") and "⚠️ 심층 분석 중 오류 발생" not in data.get("analysis", ""):
        analysis_text = data["analysis"]
        
        sec1 = extract_section(1, analysis_text)
        sec10 = extract_section(10, analysis_text)
        
        # 투자 판정 배지 추출
        decision = "분석 불가"
        badge_match = re.search(r'투자\s*(?:판정|여부|판단)\s*[:\-]?\s*\[?\s*([^\]\n\r]+)\]?', sec10, re.IGNORECASE)
        if badge_match:
            decision = badge_match.group(1).strip()
            if decision.upper() == 'NEUTRAL': decision = 'Neutral'
            elif decision.upper() == 'DANGER': decision = 'Danger'
            elif decision.upper() == 'GO': decision = 'GO'
            
        # 1. 핵심요점 & 배지
        doc.add_heading('핵심요점', level=1)
        
        p_badge = doc.add_paragraph()
        run_badge = p_badge.add_run(decision)
        run_badge.bold = True
        run_badge.font.size = Pt(14)
        if decision == 'GO' or '낙찰' in decision:
            run_badge.font.color.rgb = RGBColor(46, 160, 67) # Green
        elif decision == 'Neutral' or '변경' in decision or '대기' in decision:
            run_badge.font.color.rgb = RGBColor(210, 153, 34) # Yellow/Orange
        else:
            run_badge.font.color.rgb = RGBColor(248, 81, 73) # Red
            
        # 요약 텍스트 추가
        if sec1:
            doc.add_paragraph(sec1)
        
        doc.add_paragraph(f"\n투자 여부: [{decision}]\n")
        
        # 2. 사진부
        h_photo = doc.add_heading(level=1)
        r_photo = h_photo.add_run('사진부 (전경 및 구조)')
        r_photo.bold = True
        r_photo.font.size = Pt(18)
        try:
            if os.path.exists(img_path_photo):
                doc.add_picture(img_path_photo, width=Inches(4.5))
            if os.path.exists(img_path_map):
                doc.add_picture(img_path_map, width=Inches(4.5))
            if os.path.exists(img_path_structure):
                doc.add_picture(img_path_structure, width=Inches(4.5))
        except Exception as e:
            print(f"워드 파일 사진 삽입 오류: {e}")
            
        # 3. 각 섹션 매핑하여 작성 (2~10 -> 1~9)
        section_titles = {
            2: "I. 기본정보",
            3: "II. 물리적 현황",
            4: "III. 권리 분석",
            5: "IV. 시세 및 수익성 현황",
            6: "V. 입지 분석",
            7: "VI. 추천 입찰가",
            8: "VII. 출구 전략",
            9: "VIII. 대출 절세 전략",
            10: "IX. 최종 결론"
        }
        
        # JSON 파싱
        chart_data = None
        json_match = re.search(r'(\{[\s\S]*?"chart_data"[\s\S]*?\})', analysis_text, re.IGNORECASE)
        if json_match:
            try:
                import json
                parsed = json.loads(json_match.group(1))
                chart_data = parsed.get("chart_data")
            except Exception as e:
                print(f"차트 데이터 파싱 오류: {e}")

        def format_man(num):
            if not num: return "0"
            try:
                if isinstance(num, str):
                    num = num.replace(",", "").strip()
                val = int(float(num))
                return f"{round(val / 10000):,}만"
            except:
                return str(num)

        for ai_num, title in section_titles.items():
            if ai_num == 10 and chart_data:
                # 차트 데이터 삽입 (9. 최종결론 직전)
                doc.add_heading('수익성 시뮬레이터', level=2)
                doc.add_paragraph('낙찰가에 따른 실질(Net) 수익금 시뮬레이터', style='Intense Quote')
                doc.add_paragraph(f"시뮬레이션 입찰가\n{format_man(chart_data.get('max_bidding_price', 0))}원")
                doc.add_paragraph(f"예상 Net 수익금:\n{format_man(chart_data.get('target_profit', 0))}원")
                
                doc.add_heading('투자금 폭포수', level=2)
                doc.add_paragraph('한계 입찰가 역산 워터폴 시뮬레이션', style='Intense Quote')
                table = doc.add_table(rows=6, cols=2)
                table.style = 'Table Grid'
                rows = [
                    ('보수적 매도가', format_man(chart_data.get('conservative_sale_price', 0))),
                    ('목표 수익 (-)', format_man(chart_data.get('target_profit', 0))),
                    ('양도소득세 (-)', format_man(chart_data.get('capital_gains_tax', 0))),
                    ('CAPEX/명도 (-)', format_man(chart_data.get('capex_and_eviction', 0))),
                    ('금융/취득세 (-)', format_man(chart_data.get('interest_and_acq_tax', 0))),
                    ('한계 상한가 (=)', format_man(chart_data.get('max_bidding_price', 0)))
                ]
                for i, row in enumerate(rows):
                    table.cell(i, 0).text = row[0]
                    table.cell(i, 1).text = str(row[1])
                doc.add_paragraph("") # Add some spacing
                
                doc.add_heading('명도 및 출구전략 흐름도', level=2)
                doc.add_paragraph('타임라인별 Exit 전략 모델', style='Intense Quote')
                doc.add_paragraph(f"낙찰 및 명도\n상한가: {format_man(chart_data.get('max_bidding_price', 0))}")
                doc.add_paragraph(f"↓\n단기 매도 (1년 내)\n양도세율: {chart_data.get('short_term_tax_rate', 77)}%\n기대수익: Low")
                doc.add_paragraph(f"↓\n전월세 임대 셋팅\n보증금 회수 (투자금 방어)")
                doc.add_paragraph(f"↓\n일반 과세 매도 (2년 후)\n양도세율: {chart_data.get('long_term_tax_rate', 20)}%\n기대수익: High")
                doc.add_paragraph("")

            h = doc.add_heading(level=1)
            r = h.add_run(title)
            r.bold = True
            r.font.size = Pt(18)
            content = extract_section(ai_num, analysis_text)
            
            # JSON 블록 제거
            content = re.sub(r'```(?:json)?[\s\S]*?```', '', content, flags=re.IGNORECASE).strip()
            
            if content:
                # 마크다운 굵은 글씨(**) 처리
                lines = content.split('\n')
                for line in lines:
                    line_stripped = line.strip()
                    if not line_stripped:
                        continue
                        
                    p = doc.add_paragraph()
                    if line_stripped.startswith("- ") or line_stripped.startswith("* "):
                        p.style = 'List Bullet'
                        line_stripped = line_stripped[2:]
                        
                    parts = line_stripped.split("**")
                    for idx, part in enumerate(parts):
                        run = p.add_run(part)
                        if idx % 2 == 1:
                            run.bold = True
            else:
                doc.add_paragraph("내용이 없습니다.")
                
    else:
        doc.add_heading('에러 발생', level=1)
        doc.add_paragraph("AI 심층 분석 데이터가 유효하지 않거나 에러가 발생했습니다.")
        doc.add_paragraph(data.get("analysis", ""))

    # -------------------------------------------------------------
    # PPTX 생성을 위한 기존 파이썬 기반 휴리스틱 로직 일부 유지
    # -------------------------------------------------------------
    loan_guide = ""
    if force_no_loan:
        loan_guide = "[대출 불가 요망] 지분경매, 위반건축물, 맹지 등 은행권 대출 불가 사유 존재. 100% 현금 매수 요망."
    elif pType == '주택':
        loan_guide = f"[주택 대출] LTV: {'40% (조정지역)' if is_reg else '70% (비규제)'} / DSR: 40% 적용."
    else:
        loan_guide = f"[비주택 대출] LTV: 60~80% / RTI: {'1.5배' if pType == '상가' else '1.2배'} 이상 요구."

    strategy_msg = ""
    if has_rights:
        strategy_msg = "특수권리(대항력 등)가 포함되어 있습니다. 보수적 입찰 필수."
    elif duration == '단기':
        strategy_msg = "명확한 권리관계입니다. 명도 완료 직후 상품화 후 일반 매매로 전환 요망."
    else:
        strategy_msg = "전월세 세팅 후 최소 2년 이상 보유하여 일반 매매 전환하는 엑시트 전략 요망."

    safe_case = case_number.replace(" ", "_").replace("/", "_")
    file_path = os.path.join(output_dir, f"{safe_case}_분석보고서.docx")
    doc.save(file_path)
    
    # -------------------------------------------------------------
    # python-pptx 를 이용한 브리핑용 PPTX 자동 생성 로직 추가
    # -------------------------------------------------------------
    try:
        from pptx import Presentation
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        prs = Presentation()
        
        # Slide 1: 표지
        slide_layout = prs.slide_layouts[0] # 0: Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"부동산 경매 AI 심층 분석"
        subtitle.text = f"사건번호: {case_number}\n타겟 물건: {pType} | 대상 주소: {data.get('address', '')}"
        
        # Slide 2: 권리분석 및 리스크
        slide_layout = prs.slide_layouts[1] # 1: Title and Content
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "권리분석 및 리스크"
        tf = body_shape.text_frame
        tf.text = f"말소기준권리: {'위험 (특수권리 존재)' if has_rights else '안전 (전액 소멸)'}"
        p = tf.add_paragraph()
        p.text = f"특수 권리: {', '.join(risks) if has_rights else '발견되지 않음'}"
        p = tf.add_paragraph()
        p.text = f"대출 가이드: LTV {loan_guide.split('LTV: ')[1].split(' /')[0]} / 주택수({house_count}주택) 방어"
        
        # Slide 3: 투자 전략
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "AI 투자 전략 및 권장가"
        tf = body_shape.text_frame
        tf.text = f"지역/물건 기반 베이스라인: {int(base_rate*100)}%"
        p = tf.add_paragraph()
        p.text = f"최종 목표 입찰가: {int(target_price/100000000):.2f} 억 원"
        p = tf.add_paragraph()
        p.text = f"전략 요약: {strategy_msg}"
        
        ppt_path = os.path.join(output_dir, f"{safe_case}_브리핑자료.pptx")
        prs.save(ppt_path)
        print(f"PPTX 생성 완료: {ppt_path}")
        
        # PPTX를 PDF로 변환 (Windows comtypes 사용)
        try:
            import comtypes.client
            pdf_path = os.path.join(output_dir, f"{safe_case}_브리핑자료.pdf")
            abs_ppt_path = os.path.abspath(ppt_path)
            abs_pdf_path = os.path.abspath(pdf_path)
            
            # 파워포인트 백그라운드 실행 및 변환
            ppt_app = comtypes.client.CreateObject('PowerPoint.Application')
            presentation = ppt_app.Presentations.Open(abs_ppt_path, WithWindow=False)
            presentation.ExportAsFixedFormat(abs_pdf_path, 2) # 2 = ppFixedFormatTypePDF
            presentation.Close()
            ppt_app.Quit()
            print(f"PDF 생성 완료: {pdf_path}")
        except Exception as pdf_err:
            print(f"PDF 변환 실패: {pdf_err}")
            
    except Exception as e:
        print(f"PPTX 생성 실패: {e}")

    return file_path

def generate_analysis_doc_from_markdown(case_number, markdown_text, output_dir=DEFAULT_DOWNLOADS_DIR):
    """
    프론트엔드의 화면 내용(마크다운)을 100% 반영하여 워드 문서를 즉석 재생성합니다.
    (이미지 3장은 가로 1열 표에 나란히 배치)
    """
    os.makedirs(output_dir, exist_ok=True)
    doc = Document()
    
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Malgun Gothic'
    font.size = Pt(11)
    
    # 0. 커버 페이지 (표지)
    doc.add_heading('AI 기반 부동산 경매 심층 분석 리포트', 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph('\n')
    
    # 기존의 가로 1x3 표 이미지 삽입부는 권리분석 양식(세로 나열)으로 통합할 것이므로 표지는 아래처럼 구성
    p_cover = doc.add_paragraph()
    p_cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_cover.add_run(f"\n사건번호: {case_number}\n").bold = True
    doc.add_page_break()

    safe_case = case_number.replace(" ", "_").replace("/", "_")
    img_path_photo = os.path.join(output_dir, safe_case, "photo.jpg")
    img_path_map = os.path.join(output_dir, safe_case, "map.jpg")
    img_path_structure = os.path.join(output_dir, safe_case, "structure.jpg")

    import re
    
    def extract_section(num, text):
        pattern = r'(?:^|\n)[#\*\s]*' + str(num) + r'\.\s*[^\n]*\n([\s\S]*?)(?=(?:^|\n)[#\*\s]*\d+\.\s*|$)'
        match = re.search(pattern, text, re.IGNORECASE)
        return match.group(1).strip() if match else ""

    # JSON 파싱
    chart_data = None
    json_match = re.search(r'(\{[\s\S]*?"chart_data"[\s\S]*?\})', markdown_text, re.IGNORECASE)
    if json_match:
        try:
            import json
            parsed = json.loads(json_match.group(1))
            chart_data = parsed.get("chart_data")
        except Exception as e:
            print(f"차트 데이터 파싱 오류: {e}")

    # JSON 찌꺼기 완벽 제거
    clean_markdown = re.sub(r'```(?:json)?[\s\S]*?```', '', markdown_text, flags=re.IGNORECASE)
    
    if clean_markdown:
        sec1 = extract_section(1, clean_markdown)
        sec10 = extract_section(10, clean_markdown)
        
        # 투자 판정 배지 추출
        decision = "분석 불가"
        badge_match = re.search(r'투자\s*(?:판정|여부|판단)\s*[:\-]?\s*\[?\s*([^\]\n\r]+)\]?', sec10, re.IGNORECASE)
        if badge_match:
            decision = badge_match.group(1).strip()
            if decision.upper() == 'NEUTRAL': decision = 'Neutral'
            elif decision.upper() == 'DANGER': decision = 'Danger'
            elif decision.upper() == 'GO': decision = 'GO'
            
        # 1. 핵심요점 & 배지
        doc.add_heading('핵심요점', level=1)
        
        p_badge = doc.add_paragraph()
        run_badge = p_badge.add_run(decision)
        run_badge.bold = True
        run_badge.font.size = Pt(14)
        if decision == 'GO' or '낙찰' in decision:
            run_badge.font.color.rgb = RGBColor(46, 160, 67) # Green
        elif decision == 'Neutral' or '변경' in decision or '대기' in decision:
            run_badge.font.color.rgb = RGBColor(210, 153, 34) # Yellow/Orange
        else:
            run_badge.font.color.rgb = RGBColor(248, 81, 73) # Red
            
        # 요약 텍스트 추가
        if sec1:
            doc.add_paragraph(sec1)
        
        doc.add_paragraph(f"\n투자 여부: [{decision}]\n")
        
        # 2. 사진부
        doc.add_heading('전경사진', level=1)
        try:
            if os.path.exists(img_path_photo):
                doc.add_picture(img_path_photo, width=Inches(4.5))
            if os.path.exists(img_path_map):
                doc.add_picture(img_path_map, width=Inches(4.5))
            if os.path.exists(img_path_structure):
                doc.add_picture(img_path_structure, width=Inches(4.5))
        except Exception as e:
            print(f"워드 파일 사진 삽입 오류: {e}")
            
        # 3. 각 섹션 매핑하여 작성 (2~10 -> 1~9)
        section_titles = {
            2: "1. 기본정보",
            3: "2. 물리적 현황",
            4: "3. 권리 분석",
            5: "4. 시세 및 수익성 현황",
            6: "5. 입지 분석",
            7: "6. 추천 입찰가",
            8: "7. 출구 전략",
            9: "8. 대출 절세 전략",
            10: "9. 최종 결론"
        }
        
        def format_man(num):
            if not num: return "0"
            try:
                if isinstance(num, str):
                    num = num.replace(",", "").strip()
                val = int(float(num))
                return f"{round(val / 10000):,}만"
            except:
                return str(num)

        for ai_num, title in section_titles.items():
            if ai_num == 10 and chart_data:
                # 차트 데이터 삽입 (9. 최종결론 직전)
                doc.add_heading('수익성 시뮬레이터', level=2)
                doc.add_paragraph('낙찰가에 따른 실질(Net) 수익금 시뮬레이터', style='Intense Quote')
                doc.add_paragraph(f"시뮬레이션 입찰가\n{format_man(chart_data.get('max_bidding_price', 0))}원")
                doc.add_paragraph(f"예상 Net 수익금:\n{format_man(chart_data.get('target_profit', 0))}원")
                
                doc.add_heading('투자금 폭포수', level=2)
                doc.add_paragraph('한계 입찰가 역산 워터폴 시뮬레이션', style='Intense Quote')
                table = doc.add_table(rows=6, cols=2)
                table.style = 'Table Grid'
                rows = [
                    ('보수적 매도가', format_man(chart_data.get('conservative_sale_price', 0))),
                    ('목표 수익 (-)', format_man(chart_data.get('target_profit', 0))),
                    ('양도소득세 (-)', format_man(chart_data.get('capital_gains_tax', 0))),
                    ('CAPEX/명도 (-)', format_man(chart_data.get('capex_and_eviction', 0))),
                    ('금융/취득세 (-)', format_man(chart_data.get('interest_and_acq_tax', 0))),
                    ('한계 상한가 (=)', format_man(chart_data.get('max_bidding_price', 0)))
                ]
                for i, row in enumerate(rows):
                    table.cell(i, 0).text = row[0]
                    table.cell(i, 1).text = str(row[1])
                doc.add_paragraph("") # Add some spacing
                
                doc.add_heading('명도 및 출구전략 흐름도', level=2)
                doc.add_paragraph('타임라인별 Exit 전략 모델', style='Intense Quote')
                doc.add_paragraph(f"낙찰 및 명도\n상한가: {format_man(chart_data.get('max_bidding_price', 0))}")
                doc.add_paragraph(f"↓\n단기 매도 (1년 내)\n양도세율: {chart_data.get('short_term_tax_rate', 77)}%\n기대수익: Low")
                doc.add_paragraph(f"↓\n전월세 임대 셋팅\n보증금 회수 (투자금 방어)")
                doc.add_paragraph(f"↓\n일반 과세 매도 (2년 후)\n양도세율: {chart_data.get('long_term_tax_rate', 20)}%\n기대수익: High")
                doc.add_paragraph("")

            doc.add_heading(title, level=1)
            content = extract_section(ai_num, clean_markdown)
            
            if content:
                # 마크다운 굵은 글씨(**) 처리
                lines = content.split('\n')
                for line in lines:
                    line_stripped = line.strip()
                    if not line_stripped:
                        continue
                        
                    p = doc.add_paragraph()
                    if line_stripped.startswith("- ") or line_stripped.startswith("* "):
                        p.style = 'List Bullet'
                        line_stripped = line_stripped[2:]
                        
                    parts = line_stripped.split("**")
                    for idx, part in enumerate(parts):
                        run = p.add_run(part)
                        if idx % 2 == 1:
                            run.bold = True
            else:
                doc.add_paragraph("내용이 없습니다.")
    else:
        doc.add_paragraph("분석 데이터가 존재하지 않습니다.")

    # 안내 말씀 (면책 조항) 추가
    doc.add_paragraph("") # 빈 줄
    disclaimer_p = doc.add_paragraph()
    disclaimer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_disc = disclaimer_p.add_run("※ 안내 말씀: 본 AI 분석 보고서는 투자 의사결정을 돕기 위한 참고 자료입니다.\n실제 투자 시에는 반드시 현장 조사와 전문가 상담 등을 병행하시길 권장해 드리며,\n본 자료를 바탕으로 한 투자 결과에 대해서는 법적 책임을 지지 않음을 너그럽게 양해 부탁드립니다.")
    run_disc.font.size = Pt(10)
    run_disc.font.color.rgb = RGBColor(128, 128, 128) # 회색

    file_path = os.path.join(output_dir, f"{safe_case}_분석보고서.docx")
    doc.save(file_path)
    return file_path

if __name__ == "__main__":
    test_data = {"case_number": "성남지원 경매9계 2024타경5020", "investor_type": "매매사업자", "repair_condition": "전체수리"}
    print(generate_analysis_doc(test_data))
